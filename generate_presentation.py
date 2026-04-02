"""
Generate FYP Presentation: Multi-Agent Reinforcement Learning for
Adaptive Traffic Signal Control
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_MID    = RGBColor(0x22, 0x22, 0x3A)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)    # cyan
ACCENT2   = RGBColor(0xFF, 0x6B, 0x6B)    # coral
ACCENT3   = RGBColor(0x4E, 0xCB, 0x71)    # green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
MUTED     = RGBColor(0x99, 0x99, 0x99)
YELLOW    = RGBColor(0xFF, 0xD9, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ─────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        import lxml.etree as etree
        solid = shape.fill._fill.find(
            './/{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solid is not None:
            sr = solid[0]
            alpha_elem = etree.SubElement(
                sr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", spacing_after=Pt(6), line_spacing=1.2):
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = spacing_after
        if line_spacing:
            p.line_spacing = line_spacing


def add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(1.8),
                     width=Inches(11.5), size=18, color=LIGHT, bullet="▸"):
    box = tb(slide, left, top, width, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.line_spacing = 1.3


def title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.3), BG_MID)
    add_rect(slide, 0, Inches(1.28), W, Inches(0.04), ACCENT)
    box = tb(slide, Inches(0.7), Inches(0.15), Inches(12), Inches(0.7))
    set_text(box.text_frame, title, size=32, bold=True, color=WHITE)
    if subtitle:
        box2 = tb(slide, Inches(0.7), Inches(0.75), Inches(12), Inches(0.5))
        set_text(box2.text_frame, subtitle, size=16, color=MUTED)


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, Inches(2.5), W, Inches(2.8), BG_MID)
    add_rect(slide, 0, Inches(2.48), W, Inches(0.04), ACCENT)
    add_rect(slide, 0, Inches(5.28), W, Inches(0.04), ACCENT)
    box = tb(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2))
    set_text(box.text_frame, title, size=44, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    if subtitle:
        box2 = tb(slide, Inches(1), Inches(4.1), Inches(11), Inches(0.8))
        set_text(box2.text_frame, subtitle, size=20, color=LIGHT,
                 align=PP_ALIGN.CENTER)
    return slide


def content_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, title, subtitle)
    return slide


def add_table(slide, rows, cols, data, left, top, width, height,
              header_color=ACCENT, cell_color=LIGHT):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    for c in range(cols):
        for r in range(rows):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_MID if r == 0 else BG_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = header_color if r == 0 else cell_color
                p.font.bold = (r == 0)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


def add_page_number(slide, num, total):
    box = tb(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(box.text_frame, f"{num}/{total}", size=10, color=MUTED,
             align=PP_ALIGN.RIGHT)


# ╔═══════════════════════════════════════════════════════════════╗
# ║                    SLIDE CONTENT                             ║
# ╚═══════════════════════════════════════════════════════════════╝

slides_list = []

# ─── 1. TITLE SLIDE ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_rect(s, 0, Inches(0.5), W, Inches(0.04), ACCENT)
add_rect(s, 0, Inches(5.5), W, Inches(0.04), ACCENT)
box = tb(s, Inches(1), Inches(1.2), Inches(11.3), Inches(2))
set_text(box.text_frame,
         "Multi-Agent Reinforcement Learning\nfor Adaptive Traffic Signal Control",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.3)
box2 = tb(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1))
set_text(box2.text_frame,
         "A Parameter-Sharing DQN Approach on Real-World Urban Networks",
         size=22, color=ACCENT, align=PP_ALIGN.CENTER)
box3 = tb(s, Inches(1), Inches(5.8), Inches(11.3), Inches(1.2))
set_text(box3.text_frame,
         "Final Year Project\nDepartment of Computer Science\n2025-2026",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)
slides_list.append(s)

# ─── 2. TABLE OF CONTENTS ──────────────────────────────────────
s = content_slide("Table of Contents")
items = [
    "1.  Problem Analysis — Why Traffic Signal Control Matters",
    "2.  Existing Solutions — Traditional & Baseline Approaches",
    "3.  System Architecture — Tools, Framework & Hierarchy",
    "4.  Algorithm Design — Dueling Double DQN with MARL",
    "5.  Methodology — Training Pipeline & Reward Engineering",
    "6.  Results — Experiments, Metrics & Findings",
    "7.  Conclusion & Future Work",
]
add_bullet_slide(s, items, size=20, bullet="◆", color=WHITE)
slides_list.append(s)

# ─── 4. PROBLEM ANALYSIS ───────────────────────────────────────
s = content_slide("The Urban Traffic Challenge")
items = [
    "Urban traffic congestion costs the global economy over $1 trillion annually in lost productivity, fuel waste, and emissions",
    "Traffic signals control ~75% of urban intersections, yet most run on fixed-time plans designed decades ago",
    "Traditional signal timing cannot adapt to real-time demand fluctuations (rush hour, accidents, events)",
    "Coordinating hundreds of intersections simultaneously is a combinatorial optimization problem — intractable for classical methods",
    "Reinforcement Learning offers a data-driven approach: agents learn optimal policies through direct interaction with the environment",
    "Challenge: Scaling single-agent RL to city-wide networks with 50–200+ traffic lights requires multi-agent coordination",
]
add_bullet_slide(s, items, size=17)
slides_list.append(s)

# ─── 6. TRADITIONAL APPROACHES ─────────────────────────────────
s = content_slide("Traditional Traffic Control Methods")
items = [
    "Fixed-Time Control — Pre-computed green/red splits based on historical data; cannot adapt to real-time demand changes",
    "Actuated Control — Extends green phases based on detector input (e.g., induction loops); reactive but lacks coordination across intersections",
    "SCOOT / SCATS — Centralized adaptive systems that adjust splits and offsets using detector data; require expensive infrastructure and are proprietary",
    "Green Wave / Offset Optimization — Coordinates adjacent signals for arterial progression; only effective on linear corridors, fails in grid networks",
]
add_bullet_slide(s, items, size=17)

box = tb(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1))
set_text(box.text_frame,
         "Common Limitation: These methods optimize locally or require global knowledge — "
         "they struggle with the curse of dimensionality in large, heterogeneous networks.",
         size=15, color=YELLOW, bold=True)
slides_list.append(s)

# ─── 7. RL-BASED APPROACHES ────────────────────────────────────
s = content_slide("Reinforcement Learning Approaches")
items = [
    "Single-Agent RL (e.g., IntelliLight) — One DQN per intersection; does not scale to city-wide networks and ignores inter-intersection effects",
    "Independent Multi-Agent RL — Each intersection has its own agent; scalable but agents cannot coordinate, leading to oscillation and instability",
    "Centralized Training, Decentralized Execution (CTDE) — Training uses global state, execution is local; communication overhead is prohibitive for 50+ agents",
    "Parameter Sharing (Our Approach) — All agents share one neural network; scales to any number of intersections with constant model size",
    "Graph-based Methods (e.g., CoLight) — Use graph attention networks to model intersection topology; powerful but computationally expensive",
]
add_bullet_slide(s, items, size=17)
slides_list.append(s)

# ─── 8. BASELINE STRATEGIES ────────────────────────────────────
s = content_slide("Baseline Strategies Used in This Project")
data = [
    ["Strategy", "Mechanism", "Limitation"],
    ["Fixed-Time", "Cycles through phases with constant\ngreen durations", "Cannot adapt to real-time\ndemand fluctuations"],
    ["Max Pressure", "Selects phase with highest upstream\nqueue pressure across all candidates", "Greedy; no lookahead,\ncan oscillate rapidly"],
    ["SOTL", "Threshold-based: switches when competing\nphase demand exceeds threshold θ", "Sensitive to threshold tuning;\nno inter-intersection coordination"],
    ["Adaptive", "Pressure-aware extension with hysteresis,\nphase skipping to highest-demand phase", "Strong but hand-tuned\nparameters per network"],
    ["GA", "Offline evolutionary search over base\ngreen time and pressure scale (6 pop × 3 gen)", "Slow; cannot adapt once\ntraffic shifts mid-episode"],
]
add_table(s, 6, 3, data, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0))
slides_list.append(s)

# ─── 9. BASELINE COMPARISON CONTEXT ────────────────────────────
s = content_slide("Why Baselines Matter",
                  "Establishing a fair performance benchmark")
items = [
    "All baselines control the same subset of traffic lights under identical demand and episode length (1,200 s)",
    "SUMO's built-in actuated controller handles uncontrolled lights — this is not a trivial \"do-nothing\" default",
    "The Adaptive baseline is surprisingly strong: pressure-aware extension + phase skipping makes it hard to beat",
    "MaxPressure is the theoretical gold standard for throughput in queuing networks (Varaiya, 2013)",
    "GA provides an offline evolutionary optimisation benchmark — amortised search cost, O(1) per decision at runtime",
    "Experiment 1: 120 episodes per strategy | Experiment 2: 30 episodes, shared seed for paired comparison",
]
add_bullet_slide(s, items, size=17)
slides_list.append(s)

# ─── 11. TOOLS & FRAMEWORK ─────────────────────────────────────
s = content_slide("Technology Stack")
col1_items = [
    ("SUMO 1.25.0", "Microscopic traffic simulator with\nreal-world OSM network support"),
    ("TraCI", "Traffic Control Interface — Python API\nfor real-time signal manipulation"),
    ("PyTorch", "Deep learning framework for\nDQN training and inference"),
    ("Python 3.10+", "Core language with NumPy,\nFlask (web dashboard)"),
]
col2_items = [
    ("OpenStreetMap", "Real-world network topology\nand traffic demand data"),
    ("Flask + Leaflet.js", "Web-based visualization\ndashboard for live monitoring"),
    ("TensorBoard", "Training metric logging\nand visualization"),
    ("Git + GitHub", "Version control and\ncollaboration"),
]

for col_idx, items in enumerate([col1_items, col2_items]):
    left = Inches(0.7) if col_idx == 0 else Inches(6.8)
    for i, (name, desc) in enumerate(items):
        y = Inches(1.7) + Inches(1.35) * i
        add_rect(s, left, y, Inches(5.5), Inches(1.15), BG_MID)
        add_rect(s, left, y, Inches(0.08), Inches(1.15), ACCENT if col_idx == 0 else ACCENT3)
        bx = tb(s, left + Inches(0.3), y + Inches(0.08), Inches(5), Inches(0.4))
        set_text(bx.text_frame, name, size=17, bold=True, color=ACCENT if col_idx == 0 else ACCENT3)
        bx2 = tb(s, left + Inches(0.3), y + Inches(0.45), Inches(5), Inches(0.65))
        set_text(bx2.text_frame, desc, size=13, color=LIGHT)
slides_list.append(s)

# ─── 12. SYSTEM HIERARCHY ──────────────────────────────────────
s = content_slide("System Hierarchy",
                  "Layered architecture from simulation to decision-making")

layers = [
    ("SUMO Simulation Engine", "Microscopic vehicle movement, lane dynamics, traffic light physics", ACCENT2, Inches(1.6)),
    ("TraCI Communication Layer", "Real-time bidirectional data exchange: observations ↔ actions", YELLOW, Inches(2.8)),
    ("Traffic Simulator Wrapper", "State collection, baseline strategies, metric calculation, demand scaling", ACCENT3, Inches(4.0)),
    ("MARL Training / Evaluation", "Shared DQN agent, reward computation, experience replay, curriculum learning", ACCENT, Inches(5.2)),
]

for name, desc, color, top in layers:
    add_rect(s, Inches(1.5), top, Inches(10.3), Inches(1.0), BG_MID)
    add_rect(s, Inches(1.5), top, Inches(0.1), Inches(1.0), color)
    bx = tb(s, Inches(1.9), top + Inches(0.05), Inches(4), Inches(0.45))
    set_text(bx.text_frame, name, size=17, bold=True, color=color)
    bx2 = tb(s, Inches(1.9), top + Inches(0.45), Inches(9.5), Inches(0.5))
    set_text(bx2.text_frame, desc, size=14, color=LIGHT)

for i in range(3):
    y = Inches(2.6) + Inches(1.2) * i
    add_rect(s, Inches(6.5), y, Inches(0.04), Inches(0.25), MUTED)

slides_list.append(s)

# ─── 13. MARL PARAMETER SHARING ────────────────────────────────
s = content_slide("Multi-Agent Parameter Sharing",
                  "One network controls all traffic lights")
items = [
    "All N traffic light agents share a single DQN — model size is O(1), not O(N)",
    "Each agent receives a LOCAL observation: lane queues, speeds, vehicle counts, phase info for its own intersection",
    "Actions are executed independently per intersection — fully decentralized execution",
    "Experience from ALL intersections feeds into ONE shared replay buffer — massive data efficiency",
    "Curriculum learning: start controlling 25% of lights, gradually increase to target ratio (50%)",
    "This approach handles networks of any size without retraining — demonstrated on Cologne (79 TLS)",
]
add_bullet_slide(s, items, size=17)

add_rect(s, Inches(3), Inches(5.8), Inches(7), Inches(1.2), BG_MID)
bx = tb(s, Inches(3.3), Inches(5.9), Inches(6.5), Inches(1))
set_text(bx.text_frame,
         "Key Insight: Parameter sharing transforms a 40-agent problem into a single-agent "
         "problem with 40× more training data per episode.",
         size=15, color=YELLOW, bold=True)
slides_list.append(s)

# ─── 15. DQN FUNDAMENTALS ──────────────────────────────────────
s = content_slide("Deep Q-Network (DQN) Fundamentals")
items = [
    "Q-Learning: Learn Q(s,a) — the expected cumulative reward for taking action a in state s",
    "Deep Q-Network: Approximate Q(s,a) with a neural network parameterized by θ",
    "Target network (θ⁻) updated via hard copy every C = 2,000 gradient steps to stabilize training",
    "Experience Replay: Store transitions (s, a, r, s') in a buffer, sample mini-batches to break correlation",
    "ε-greedy exploration: Start with random actions (ε=1.0), linearly decay to ε=0.05 over training",
    "Loss function: Smooth L1 (Huber) loss between predicted Q and target Q values",
]
add_bullet_slide(s, items, size=17)

add_rect(s, Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.2), BG_MID)
bx = tb(s, Inches(2.8), Inches(5.65), Inches(7.8), Inches(1.1))
set_text(bx.text_frame,
         "Q-target = r + γⁿ · Q_target(s', argmax_a Q_online(s', a))\n"
         "Loss = SmoothL1( Q_online(s, a) − Q-target )",
         size=16, color=ACCENT, align=PP_ALIGN.CENTER, font_name="Consolas")
slides_list.append(s)

# ─── 16. ENHANCEMENTS ──────────────────────────────────────────
s = content_slide("DQN Enhancements Used")

enhancements = [
    ("Double DQN", "Action selection from online network, value from target network — eliminates overestimation bias", ACCENT),
    ("Dueling Architecture", "Separates Q into V(s) + A(s,a) streams — learns state values independently of action advantages", ACCENT3),
    ("Prioritized Experience Replay", "Samples transitions proportional to |TD error|^α — focuses learning on surprising experiences", YELLOW),
    ("N-step Returns (n=3)", "Bootstrap from 3 steps ahead: r₁ + γr₂ + γ²r₃ + γ³V(s₃) — faster credit assignment", ACCENT2),
    ("Action Masking", "Prevents invalid switches before min_green_seconds (10s) — reduces wasted exploration", WHITE),
]

for i, (name, desc, color) in enumerate(enhancements):
    y = Inches(1.6) + Inches(1.05) * i
    add_rect(s, Inches(0.7), y, Inches(11.9), Inches(0.9), BG_MID)
    add_rect(s, Inches(0.7), y, Inches(0.08), Inches(0.9), color)
    bx = tb(s, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.4))
    set_text(bx.text_frame, name, size=16, bold=True, color=color)
    bx2 = tb(s, Inches(1.0), y + Inches(0.4), Inches(11.3), Inches(0.45))
    set_text(bx2.text_frame, desc, size=13, color=LIGHT)
slides_list.append(s)

# ─── 17. NETWORK ARCHITECTURE ──────────────────────────────────
s = content_slide("Neural Network Architecture",
                  "Dueling DQN with LayerNorm")

arch_layers = [
    ("Input Layer", f"State vector: 12 lanes × 3 features + 16-dim phase one-hot + time + 4 pressure features = 57 dims", Inches(1.6)),
    ("Hidden Layer 1", "256 neurons → LayerNorm → ReLU", Inches(2.5)),
    ("Hidden Layer 2", "256 neurons → LayerNorm → ReLU", Inches(3.2)),
    ("Hidden Layer 3", "128 neurons → LayerNorm → ReLU", Inches(3.9)),
    ("Value Stream", "128 → 1   (state value V(s))", Inches(4.8)),
    ("Advantage Stream", "128 → 2   (action advantages A(s,a))", Inches(5.5)),
    ("Output", "Q(s,a) = V(s) + A(s,a) − mean(A)   →   2 actions: KEEP / SWITCH", Inches(6.3)),
]

for name, desc, top in arch_layers:
    w = Inches(10) if "Output" not in name else Inches(10)
    add_rect(s, Inches(1.6), top, Inches(10.1), Inches(0.65), BG_MID)
    bx = tb(s, Inches(1.9), top + Inches(0.02), Inches(3), Inches(0.35))
    set_text(bx.text_frame, name, size=14, bold=True, color=ACCENT)
    bx2 = tb(s, Inches(4.5), top + Inches(0.02), Inches(7), Inches(0.55))
    set_text(bx2.text_frame, desc, size=13, color=LIGHT)
slides_list.append(s)

# ─── 19. STATE & ACTION SPACE ──────────────────────────────────
s = content_slide("MDP Formulation",
                  "State space, action space, and reward design")
box = tb(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5))
tf = box.text_frame
tf.word_wrap = True

sections = [
    ("State Space (per intersection):", [
        "Vehicle counts per lane (normalized by 20)",
        "Queue lengths per lane (normalized by 15)",
        "Mean speeds per lane (normalized by 20)",
        "Current phase (16-dim one-hot encoding)",
        "Remaining green time (normalized by 60s)",
        "Pressure features: current/next pressure,",
        "  pressure delta, total queue (4 dims)",
    ]),
    ("Action Space:", [
        "Action 0: KEEP current phase",
        "Action 1: SWITCH to next green phase",
        "Action masking: block switches before",
        "  minimum green time (10 seconds)",
    ]),
]

for sec_title, bullets in sections:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = sec_title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    p.space_after = Pt(4)
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"  ▸  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(8)

add_rect(s, Inches(7), Inches(1.6), Inches(5.8), Inches(3.5), BG_MID)
bx = tb(s, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.4))
set_text(bx.text_frame, "Reward Function:", size=17, bold=True, color=ACCENT)
reward_lines = [
    ("r_wait", "= −0.50 × clip(total_wait / scale, 0, 1)", "Penalize waiting"),
    ("r_queue", "= +0.35 × clip(Δqueue / scale, −1, 1)", "Reward queue reduction"),
    ("r_switch", "= −0.01  (if action = SWITCH)", "Small switch cost"),
    ("R(s,a)", "= r_wait + r_queue + r_switch", "Total reward"),
]
for i, (name, formula, desc) in enumerate(reward_lines):
    y = Inches(2.3) + Inches(0.7) * i
    bx = tb(s, Inches(7.3), y, Inches(1.5), Inches(0.35))
    set_text(bx.text_frame, name, size=14, bold=True, color=YELLOW, font_name="Consolas")
    bx2 = tb(s, Inches(8.6), y, Inches(3), Inches(0.35))
    set_text(bx2.text_frame, formula, size=12, color=WHITE, font_name="Consolas")
    bx3 = tb(s, Inches(8.6), y + Inches(0.3), Inches(3.5), Inches(0.3))
    set_text(bx3.text_frame, desc, size=11, color=MUTED)
slides_list.append(s)

# ─── 20. TRAINING PIPELINE ─────────────────────────────────────
s = content_slide("Training Pipeline",
                  "Episode-based training with curriculum learning")

steps = [
    ("1", "Initialize", "Load OSM network, create SUMO instance,\nselect controlled TLS subset"),
    ("2", "Curriculum", "Start at 25% control ratio, linearly\nramp to 50% over first 100 episodes"),
    ("3", "Simulate", "Step SUMO, collect per-TLS observations\nevery decision-interval (5-10s)"),
    ("4", "Act", "ε-greedy action selection with masking;\napply setPhaseDuration via TraCI"),
    ("5", "Learn", "Store in PER buffer, sample batch of 128,\ncompute n-step TD targets, backprop"),
    ("6", "Adapt", "Cosine LR schedule, plateau recovery\n(boost ε if reward stagnates)"),
]

for i, (num, name, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    left = Inches(0.7) + Inches(4.1) * col
    top = Inches(1.6) + Inches(2.7) * row
    add_rect(s, left, top, Inches(3.8), Inches(2.3), BG_MID)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    cbox = tb(s, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5))
    set_text(cbox.text_frame, num, size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    bx = tb(s, left + Inches(0.8), top + Inches(0.15), Inches(2.8), Inches(0.4))
    set_text(bx.text_frame, name, size=16, bold=True, color=ACCENT)
    bx2 = tb(s, left + Inches(0.2), top + Inches(0.7), Inches(3.4), Inches(1.5))
    set_text(bx2.text_frame, desc, size=13, color=LIGHT)
slides_list.append(s)

# ─── 21. KEY HYPERPARAMETERS ────────────────────────────────────
s = content_slide("Key Hyperparameters")
data = [
    ["Parameter", "Value", "Notes"],
    ["Learning Rate", "5×10⁻⁵ → 1×10⁻⁵", "Cosine annealing schedule"],
    ["Discount Factor (γ)", "0.97", "High for long-horizon planning"],
    ["N-step Returns", "n = 3", "Faster credit assignment"],
    ["Batch Size", "128", "Mini-batch from PER buffer"],
    ["Replay Buffer", "200,000", "PER: α=0.6, β: 0.4→1.0"],
    ["Target Update", "Hard copy every C=2,000", "Prevents moving-target instability"],
    ["ε Schedule", "1.0 → 0.05 (linear)", "Plateau recovery boosts ε if stalled"],
    ["Hidden Layers", "[256, 256, 128] + LayerNorm", "Dueling streams after layer 3"],
    ["Decision Interval", "Exp 1: 5 s | Exp 2 MARL: 10 s", "Longer step cuts yellow share of RL window"],
    ["Min Green Time", "10 seconds", "Action mask enforced"],
    ["Curriculum", "25% → 50% over 100 eps", "Gradual intersection ramp"],
]
add_table(s, 12, 3, data, Inches(1), Inches(1.5), Inches(11.3), Inches(5.7))
slides_list.append(s)

# ─── 23. DATASET ───────────────────────────────────────────────
s = content_slide("Dataset: Cologne, Germany (OpenStreetMap)",
                  "Real-world urban network from the Cologne scenario")

stats_left = [
    ("Road Network", "194 km total road length"),
    ("Edges / Lanes", "3,070 edges, 5,618 lanes"),
    ("Junctions", "2,512 total (79 signalized)"),
    ("Avg Speed Limit", "11.2 m/s (40 km/h)"),
]
stats_right = [
    ("Passenger Trips", "9,537 vehicles in scenario"),
    ("Truck Trips", "638 heavy vehicles"),
    ("Bus Routes", "41 public transport lines"),
    ("Pedestrians", "4,402 pedestrian trips"),
]

for col_idx, items in enumerate([stats_left, stats_right]):
    left = Inches(0.7) if col_idx == 0 else Inches(6.8)
    for i, (name, val) in enumerate(items):
        y = Inches(1.7) + Inches(1.2) * i
        add_rect(s, left, y, Inches(5.5), Inches(1.0), BG_MID)
        add_rect(s, left, y, Inches(0.08), Inches(1.0), ACCENT if col_idx == 0 else ACCENT3)
        bx = tb(s, left + Inches(0.3), y + Inches(0.08), Inches(5), Inches(0.4))
        set_text(bx.text_frame, name, size=16, bold=True, color=ACCENT if col_idx == 0 else ACCENT3)
        bx2 = tb(s, left + Inches(0.3), y + Inches(0.45), Inches(5), Inches(0.5))
        set_text(bx2.text_frame, val, size=14, color=LIGHT)

box = tb(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.8))
set_text(box.text_frame,
         "Episode duration: 1,200 s (20 min) | Demand: scale 0.7 (Exp 1) & 0.85 (Exp 2) | "
         "Evaluation: 120 episodes (Exp 1), 30 paired-seed episodes (Exp 2)",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)
slides_list.append(s)

# ─── 24. EXPERIMENT DESIGN ─────────────────────────────────────
s = content_slide("Experiment Design",
                  "Two experiments probing different operating regimes")
items = [
    "Experiment 1 — Moderate congestion: scale 0.7, ratio 0.5 (40/79 TLS), Δt = 5 s for all, 120 episodes, MARL v8 (K=20 lane padding)",
    "Experiment 2 — Elevated congestion: scale 0.85, ratio 0.2 (16/79 TLS), 30 paired-seed episodes, MARL v16 (K=12, c_switch=0.01)",
    "Exp 2 asymmetry: MARL uses Δt = 10 s so yellow clearance is ~30% of each window (vs ~60% at 5 s); baselines stay at 5 s",
    "Baselines: Fixed-Time, MaxPressure, SOTL, Adaptive, GA — all under identical conditions per experiment",
    "Traffic metrics: avg waiting time, throughput (veh/hr), congestion index, avg queue length, avg speed",
    "Environmental metrics: CO₂, NOₓ, PMₓ, fuel — plus CO₂ per arrived vehicle (Experiment 1)",
    "Statistical analysis: mean ± std, 95% CI, Welch's t-test (Exp 1) / paired t-test (Exp 2)",
]
add_bullet_slide(s, items, size=16)
slides_list.append(s)

# ─── 25. RESULTS: EXPERIMENT 1 ─────────────────────────────────
s = content_slide("Experiment 1: Moderate Congestion",
                  "Scale 0.7 | 50% control (40/79 TLS) | Δt = 5 s | 120 episodes")
data = [
    ["Strategy", "Avg Wait (s)", "Throughput (veh/hr)", "Congest. Idx", "Avg Speed (m/s)"],
    ["SOTL", "94.80 ± 1.20", "759.0 ± 40.3", "0.72 ± 0.02", "0.63 ± 0.01"],
    ["Adaptive", "95.26 ± 1.20", "762.2 ± 41.2", "0.69 ± 0.02", "0.63 ± 0.01"],
    ["MARL DQN ★", "95.34 ± 1.30", "775.8 ± 40.1", "0.64 ± 0.02", "0.63 ± 0.01"],
    ["MaxPressure", "95.68 ± 1.25", "773.9 ± 32.9", "0.66 ± 0.02", "0.63 ± 0.01"],
    ["Fixed-Time", "95.93 ± 1.21", "692.4 ± 35.6", "0.74 ± 0.01", "0.61 ± 0.01"],
]
tbl = add_table(s, 6, 5, data, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.0))

box = tb(s, Inches(0.6), Inches(6.0), Inches(12), Inches(1.2))
set_text(box.text_frame,
         "MARL leads: +12.0% throughput, −13.7% congestion, +29.6% throughput efficiency vs Fixed-Time.\n"
         "Waiting time within 0.5 s of best (SOTL) — the 50% uncontrolled TLS dominate the network average.",
         size=15, color=ACCENT3, bold=True)
slides_list.append(s)

# ─── 26. RESULTS: EXPERIMENT 2 ─────────────────────────────────
s = content_slide("Experiment 2: Elevated Congestion",
                  "Scale 0.85 | 20% control (16/79 TLS) | Δt = 10 s MARL, 5 s baselines | 30 paired episodes")
data = [
    ["Strategy", "Avg Wait (s)", "Throughput (veh/ep)", "Avg Speed (m/s)", "Avg Queue"],
    ["Adaptive", "100.00 ± 1.00", "219.9 ± 14.0", "0.538 ± 0.008", "2660 ± 14"],
    ["MARL DQN ★", "100.11 ± 0.96", "218.6 ± 11.1", "0.534 ± 0.008", "2665 ± 10"],
    ["MaxPressure", "100.23 ± 1.08", "218.2 ± 12.5", "0.537 ± 0.009", "2662 ± 13"],
    ["GA", "100.54 ± 1.12", "212.9 ± 10.7", "0.532 ± 0.007", "2670 ± 9"],
    ["Fixed-Time", "100.75 ± 1.03", "213.0 ± 11.6", "0.534 ± 0.008", "2670 ± 9"],
    ["SOTL", "100.86 ± 1.37", "213.8 ± 12.7", "0.533 ± 0.010", "2670 ± 12"],
]
tbl = add_table(s, 7, 5, data, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.5))

box = tb(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.8))
set_text(box.text_frame,
         "MARL places 2nd — only 0.11 s behind Adaptive — despite being trained at ratio 0.5 and evaluated at 0.2.\n"
         "Full spread under 1 s: only 16/79 TLS controlled, leaving little leverage for any strategy.",
         size=15, color=YELLOW, bold=True)
slides_list.append(s)

# ─── 26b. ENVIRONMENTAL RESULTS ────────────────────────────────
s = content_slide("Environmental Impact",
                  "Emissions barely differ — but per-vehicle CO₂ tells a different story")
items = [
    "Total CO₂ nearly identical across all strategies: spread < 0.2% in both experiments",
    "Dominant factor is aggregate traffic volume and SUMO's HBEFA speed–acceleration model, not signal timing",
    "The interesting metric is CO₂ per arrived vehicle (Experiment 1):",
    "    MARL & MaxPressure: 2.42 × 10⁷ mg — a 10.7% drop vs Fixed-Time (2.71 × 10⁷ mg)",
    "    This is a throughput effect: higher throughput spreads the same emission budget over more completed trips",
    "Signal control determines WHERE vehicles idle, not HOW MUCH they emit in total",
    "Exp 2 even more compressed: total CO₂ ranges 7.110–7.122 × 10⁹ mg across all six strategies",
]
add_bullet_slide(s, items, size=16)
slides_list.append(s)

# ─── 27. ANALYSIS: KEY FINDINGS ────────────────────────────────
s = content_slide("Key Findings & Analysis")
items = [
    "MARL learns a conservative policy: holds current phase unless pressure gap clearly justifies the yellow-time cost of switching",
    "Experiment 1: MARL leads on throughput (+12.0%), congestion (−13.7%), and throughput efficiency (+29.6%) over Fixed-Time",
    "Experiment 2: MARL places 2nd (0.11 s behind Adaptive) despite being trained at ratio 0.5 and evaluated at 0.2 — robust to distribution shift",
    "Switch rate stabilises at ~15%, up from ~25% early in training — the agent identifies specific states where switching pays off",
    "Parameter sharing: one small MLP controls 16–40 intersections with zero additional parameters",
    "Adaptive baseline deserves credit — consistently 1st or 2nd, but requires hand-tuned parameters per network",
    "Emissions: per-arrived-vehicle CO₂ drops 10.7% vs Fixed-Time (throughput effect, not cleaner driving)",
]
add_bullet_slide(s, items, size=16)
slides_list.append(s)

# ─── 28. CHALLENGES & LESSONS ──────────────────────────────────
s = content_slide("Challenges & Lessons Learned")
items = [
    "Δt vs yellow clearance: at 5 s steps, ~3 s yellow is most of the window—RL saw little benefit from switching; 10 s steps restored learning",
    "Stabilisation suite required: removing any one of {hard target copy, Q-target clamp, gradient clip, LayerNorm} caused Q-value divergence",
    "Reward engineering: iterated from 7-term to 3-term reward; simpler is better — dominant terms are waiting-time penalty and queue-reduction bonus",
    "Partial control (20–50%) compresses all network-wide metrics — makes differences hard to see but reflects realistic deployment constraints",
    "Unpaired seeds in Exp 1 (different base seed per strategy) weaken the comparison — Exp 2 fixes this with shared seed 500",
    "Curriculum learning was essential: training on all 40 TLS from episode 1 produced noisy gradients and no convergence",
]
add_bullet_slide(s, items, size=16)
slides_list.append(s)

# ─── 29. FUTURE WORK ───────────────────────────────────────────
s = content_slide("Future Work")
items = [
    "Full network control: Scale the controlled fraction toward 100% and measure whether MARL's advantage grows",
    "Cross-city transfer: Test the Cologne-trained policy on Berlin, Los Angeles — if it fails, explore fine-tuning or meta-learning",
    "Inter-agent communication: Add a message-passing mechanism (e.g., attention over neighbours) for corridor-level coordination",
    "Policy-gradient methods: Replace DQN with PPO or SAC to enable continuous action spaces (directly setting green durations)",
    "Richer action space: Allow direct phase selection (5–8 actions) instead of binary keep/switch for more targeted control",
    "Sim-to-real: Bridge the gap through domain randomisation or calibration against real detector counts",
]
add_bullet_slide(s, items, size=16)
slides_list.append(s)

# ─── 30. CONCLUSION ────────────────────────────────────────────
s = content_slide("Conclusion")
items = [
    "Reproducible MARL benchmarking framework: OSM import → SUMO simulation → DQN training → multi-strategy evaluation with emissions",
    "Scalability under heterogeneous geometry: one DQN controls 16–40 intersections differing in lanes, phases, and connectivity",
    "Comprehensive evaluation: 5 baselines + GA, 10 metrics, 2 congestion regimes, traffic + environmental analysis",
    "Exp 1: +12% throughput, −13.7% congestion, +29.6% throughput efficiency vs Fixed-Time (120 episodes)",
    "Exp 2: 2nd place (0.11 s behind Adaptive) despite training/eval distribution shift and Δt asymmetry (30 paired episodes)",
    "Interface lesson: match RL decision interval to signal yellow/all-red in SUMO—short Δt can dominate what the agent can learn",
]
add_bullet_slide(s, items, size=17)

add_rect(s, Inches(2), Inches(5.8), Inches(9.3), Inches(1.2), BG_MID)
bx = tb(s, Inches(2.3), Inches(5.9), Inches(8.8), Inches(1))
set_text(bx.text_frame,
         "\"The agent learns not just when to switch signals, but crucially, when not to —\n"
         "discovering that the cost of action often outweighs its benefit.\"",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
slides_list.append(s)

# ─── 31. THANK YOU ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_rect(s, 0, Inches(2.5), W, Inches(3), BG_MID)
add_rect(s, 0, Inches(2.48), W, Inches(0.04), ACCENT)
add_rect(s, 0, Inches(5.48), W, Inches(0.04), ACCENT)
box = tb(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1))
set_text(box.text_frame, "Thank You", size=48, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
box2 = tb(s, Inches(1), Inches(4.0), Inches(11.3), Inches(1))
set_text(box2.text_frame, "Questions & Discussion",
         size=24, color=LIGHT, align=PP_ALIGN.CENTER)
slides_list.append(s)

# ── Add page numbers ───────────────────────────────────────────
total = len(prs.slides)
for i, slide in enumerate(prs.slides):
    add_page_number(slide, i + 1, total)

# ── Save ───────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__),
                        "FYP_Presentation_MARL_Traffic_Signal_Control.pptx")
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Total slides: {total}")
